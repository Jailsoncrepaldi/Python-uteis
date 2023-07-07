import requests
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def buscar_letra_musica(artist, song):
    url = f"https://www.letras.mus.br/{artist.replace(' ', '-').lower()}/{song.replace(' ', '-').lower()}/"
    response = requests.get(url)
    soup = BeautifulSoup(response.text, "html.parser")
    letra_div = soup.find("div", class_="cnt-letra")
    if letra_div:
        letra = letra_div.get_text(separator="\n").strip()
        return letra, artist, song
    else:
        return None


def criar_slides(letra, artist, song):
    presentation = Presentation()
    lines = letra.split("\n")

    # Slide do título da música
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    title.text = song

    # Definir cor de fundo dos slides
    for slide in presentation.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Cor de fundo preta

    # Contador para controlar a cada 4 linhas
    counter = 0
    slide_content = ""

    for line in lines:
        if counter < 4:
            slide_content += line + "\n"
            counter += 1
        else:
            slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            content = slide.placeholders[0]
            content.text = slide_content.strip()

            # Remove as bolinhas nas linhas
            for paragraph in content.text_frame.paragraphs:
                paragraph.level = 0

            slide_content = line + "\n"
            counter = 1

    # Adiciona o último slide
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    content = slide.placeholders[0]
    content.text = slide_content.strip()

    # Remove as bolinhas nas linhas
    for paragraph in content.text_frame.paragraphs:
        paragraph.level = 0

    # Ajustar a cor da fonte
    for slide in presentation.slides:
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Cor de fundo preta
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor(255, 255, 255)  # Cor do texto (branco)

    return presentation


def main():
    artist = input("Digite o nome do artista: ")
    song = input("Digite o nome da música: ")

    letra = buscar_letra_musica(artist, song)

    if letra:
        presentation = criar_slides(letra[0], letra[1], letra[2])
        presentation.save(f"{letra[1]}_{letra[2]}_letra.pptx")
        print(f"Slides gerados com sucesso no arquivo: {letra[1]}_{letra[2]}_letra.pptx")
    else:
        print("Não foi possível encontrar a letra da música.")
        input()


if __name__ == "__main__":
    main()
